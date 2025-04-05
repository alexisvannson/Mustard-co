import streamlit as st
import os
import tempfile
import io
from PyPDF2 import PdfReader
from pptx import Presentation
import requests
from bs4 import BeautifulSoup
import re
import nltk
from nltk.tokenize import sent_tokenize
import numpy as np
import faiss
import boto3
from langchain_aws import BedrockEmbeddings
from langchain_openai import OpenAIEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
import tiktoken

# Download NLTK data
nltk.download('punkt', quiet=True)

# Set up the page
st.set_page_config(page_title="Document Parser & Vector Search", layout="wide")
st.title("Document Parser & Vector Search")

# Initialize session state variables
if 'chunks' not in st.session_state:
    st.session_state.chunks = []
if 'faiss_index' not in st.session_state:
    st.session_state.faiss_index = None
if 'embeddings' not in st.session_state:
    st.session_state.embeddings = None
if 'document_content' not in st.session_state:
    st.session_state.document_content = ""
if 'parsed_documents' not in st.session_state:
    st.session_state.parsed_documents = []

# Functions for document parsing
def extract_text_from_pdf(pdf_file):
    reader = PdfReader(pdf_file)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text

def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_html(html_content):
    soup = BeautifulSoup(html_content, 'html.parser')
    
    # Remove script and style elements
    for script in soup(["script", "style"]):
        script.extract()
    
    # Extract text
    text = soup.get_text()
    
    # Clean up text
    lines = (line.strip() for line in text.splitlines())
    chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
    text = '\n'.join(chunk for chunk in chunks if chunk)
    
    return text

def extract_text_from_notion(notion_url):
    # This is a simplified version - real implementation would use Notion API
    response = requests.get(notion_url)
    if response.status_code == 200:
        return extract_text_from_html(response.text)
    else:
        return f"Failed to fetch Notion page: {response.status_code}"

def extract_text_from_youtube(youtube_url):
    # Simplified implementation - real version would use YouTube transcript API
    st.warning("YouTube transcript extraction is a placeholder. In a real app, this would use the YouTube API.")
    return f"Placeholder text for transcript from: {youtube_url}"

def estimate_token_count(text, encoding_name="cl100k_base"):
    """Estimate the number of tokens in a text using tiktoken."""
    encoding = tiktoken.get_encoding(encoding_name)
    return len(encoding.encode(text))

def smart_chunk_text(text, chunk_size=512, chunk_overlap=100):
    """
    Intelligently chunk text respecting semantic boundaries like paragraphs and sentences.
    """
    # Initialize text splitter
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=estimate_token_count,
        separators=["\n\n", "\n", ". ", " ", ""]
    )
    
    chunks = text_splitter.split_text(text)
    
    # Further process chunks to ensure they respect sentence boundaries when possible
    refined_chunks = []
    for chunk in chunks:
        # Check if chunk ends mid-sentence
        if chunk.endswith(".") or chunk.endswith("?") or chunk.endswith("!"):
            refined_chunks.append(chunk)
        else:
            # Try to find a good breaking point
            sentences = sent_tokenize(chunk)
            if len(sentences) > 1:
                # Join all sentences except the last one
                refined_chunks.append(" ".join(sentences[:-1]))
                # Store the last incomplete sentence with the next chunk
                if len(chunks) > chunks.index(chunk) + 1:
                    chunks[chunks.index(chunk) + 1] = sentences[-1] + " " + chunks[chunks.index(chunk) + 1]
            else:
                refined_chunks.append(chunk)
                
    return refined_chunks

# Setup sidebar configuration
with st.sidebar:
    st.header("Settings")
    
    # Embedding provider selection
    embedding_provider = st.selectbox(
        "Embedding Provider",
        ["Amazon Bedrock", "OpenAI", "WAS VM (NeMo)", "Mistral API"]
    )
    
    # Provider-specific settings
    if embedding_provider == "Amazon Bedrock":
        aws_region = st.text_input("AWS Region", "us-east-1")
        embedding_model = st.selectbox(
            "Embedding Model", 
            ["amazon.titan-embed-text-v1", "cohere.embed-english-v3", "cohere.embed-multilingual-v3"]
        )
    elif embedding_provider == "OpenAI":
        api_key = st.text_input("OpenAI API Key", type="password")
        os.environ["OPENAI_API_KEY"] = api_key
    elif embedding_provider == "WAS VM (NeMo)":
        nemo_url = st.text_input("NeMo API Endpoint", "http://localhost:8000/embeddings")
    elif embedding_provider == "Mistral API":
        mistral_api_key = st.text_input("Mistral API Key", type="password")
    
    st.header("Chunking Settings")
    chunk_size = st.slider("Chunk Size (tokens)", 128, 1024, 512)
    chunk_overlap = st.slider("Chunk Overlap (tokens)", 0, 200, 100)

# Get embedding model based on selected provider
def get_embeddings_model():
    if embedding_provider == "Amazon Bedrock":
        bedrock_client = boto3.client(
            service_name="bedrock-runtime",
            region_name=aws_region,
        )
        
        return BedrockEmbeddings(
            client=bedrock_client,
            model_id=embedding_model
        )
    elif embedding_provider == "OpenAI":
        if not api_key:
            st.error("Please enter an OpenAI API key.")
            return None
        return OpenAIEmbeddings()
    elif embedding_provider == "WAS VM (NeMo)":
        # This is a placeholder for NeMo embeddings integration
        # In a real implementation, you would use a custom embeddings class
        class NemoEmbeddings:
            def __init__(self, api_url):
                self.api_url = api_url
            
            def embed_query(self, text):
                # This would make an API call to your NeMo instance
                # Placeholder implementation
                response = requests.post(self.api_url, json={"text": text})
                if response.status_code == 200:
                    return response.json()["embedding"]
                else:
                    raise Exception(f"Failed to get embeddings: {response.status_code}")
                    
            def embed_documents(self, documents):
                # Embed multiple documents
                return [self.embed_query(doc) for doc in documents]
                
        return NemoEmbeddings(nemo_url)
    elif embedding_provider == "Mistral API":
        # Placeholder for Mistral API embeddings
        # In a real implementation, you would use their API
        class MistralEmbeddings:
            def __init__(self, api_key):
                self.api_key = api_key
                self.api_url = "https://api.mistral.ai/v1/embeddings"
                
            def embed_query(self, text):
                headers = {
                    "Authorization": f"Bearer {self.api_key}",
                    "Content-Type": "application/json"
                }
                response = requests.post(
                    self.api_url,
                    headers=headers,
                    json={"input": text, "model": "mistral-embed"}
                )
                if response.status_code == 200:
                    return response.json()["data"][0]["embedding"]
                else:
                    raise Exception(f"Failed to get embeddings: {response.status_code}")
                    
            def embed_documents(self, documents):
                return [self.embed_query(doc) for doc in documents]
                
        return MistralEmbeddings(mistral_api_key)

# File uploader
uploaded_file = st.file_uploader("Upload a document (PDF, PPTX, HTML, TXT)", type=["pdf", "pptx", "html", "txt"])

# URL input
url_input = st.text_input("Or enter a URL (Website, Notion, YouTube)")

# Parse document button
if st.button("Parse Document"):
    if uploaded_file is not None:
        file_extension = uploaded_file.name.split(".")[-1].lower()
        
        with st.spinner("Parsing document..."):
            if file_extension == "pdf":
                st.session_state.document_content = extract_text_from_pdf(uploaded_file)
            elif file_extension == "pptx":
                st.session_state.document_content = extract_text_from_pptx(uploaded_file)
            elif file_extension == "html":
                content = uploaded_file.read().decode('utf-8')
                st.session_state.document_content = extract_text_from_html(content)
            elif file_extension == "txt":
                st.session_state.document_content = uploaded_file.read().decode('utf-8')
                
        st.success(f"Document parsed successfully! Content length: {len(st.session_state.document_content)} characters")
        
        # Display sample of parsed content
        st.subheader("Sample of Parsed Content")
        st.text_area("", st.session_state.document_content[:1000] + "...", height=200)
        
        # Save doc info
        doc_info = {
            "name": uploaded_file.name,
            "type": file_extension,
            "size": len(st.session_state.document_content)
        }
        st.session_state.parsed_documents.append(doc_info)
        
    elif url_input:
        with st.spinner("Parsing content from URL..."):
            if "notion.so" in url_input:
                st.session_state.document_content = extract_text_from_notion(url_input)
            elif "youtube.com" in url_input or "youtu.be" in url_input:
                st.session_state.document_content = extract_text_from_youtube(url_input)
            else:
                # Treat as general website
                try:
                    response = requests.get(url_input)
                    st.session_state.document_content = extract_text_from_html(response.text)
                except Exception as e:
                    st.error(f"Failed to fetch URL: {e}")
                    
        st.success(f"Content parsed successfully! Content length: {len(st.session_state.document_content)} characters")
        
        # Display sample of parsed content
        st.subheader("Sample of Parsed Content")
        st.text_area("", st.session_state.document_content[:1000] + "...", height=200)
        
        # Save doc info
        doc_info = {
            "name": url_input,
            "type": "url",
            "size": len(st.session_state.document_content)
        }
        st.session_state.parsed_documents.append(doc_info)
    else:
        st.warning("Please upload a file or enter a URL to parse.")

# Process into chunks and create embeddings
if st.session_state.document_content and st.button("Process & Create Embeddings"):
    # Check if embedding provider is properly configured
    embeddings_model = get_embeddings_model()
    
    if embeddings_model is None:
        st.error("Please configure your embedding provider settings properly.")
    else:
        with st.spinner("Chunking text and creating embeddings..."):
            # Smart chunking
            st.session_state.chunks = smart_chunk_text(
                st.session_state.document_content, 
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap
            )
            
            # Create embeddings for each chunk
            chunk_embeddings = []
            
            # Process in batches to avoid rate limits and improve performance
            batch_size = 10
            for i in range(0, len(st.session_state.chunks), batch_size):
                batch = st.session_state.chunks[i:i+batch_size]
                
                # Update progress
                progress_text = f"Processing chunks {i+1} to {min(i+batch_size, len(st.session_state.chunks))} of {len(st.session_state.chunks)}"
                progress_bar = st.progress(i / len(st.session_state.chunks))
                st.text(progress_text)
                
                # Get embeddings for the batch
                batch_embeddings = embeddings_model.embed_documents(batch)
                chunk_embeddings.extend(batch_embeddings)
                
                # Update progress
                progress_bar.progress(min((i + batch_size) / len(st.session_state.chunks), 1.0))
            
            # Create FAISS index
            dimension = len(chunk_embeddings[0])
            index = faiss.IndexFlatL2(dimension)
            index.add(np.array(chunk_embeddings).astype('float32'))
            
            st.session_state.faiss_index = index
            st.session_state.embeddings = embeddings_model
            
        st.success(f"Created {len(st.session_state.chunks)} chunks and built vector index!")
        
        # Display chunk info
        st.subheader("Chunk Information")
        for i, chunk in enumerate(st.session_state.chunks[:5]):
            st.text_area(f"Chunk {i+1}", chunk[:200] + "..." if len(chunk) > 200 else chunk, height=100)
        
        if len(st.session_state.chunks) > 5:
            st.write(f"... and {len(st.session_state.chunks) - 5} more chunks")

# Ask a question section
st.header("Ask a Question")
query = st.text_input("Enter your question:")

if query and st.button("Search") and st.session_state.faiss_index is not None:
    with st.spinner("Searching for relevant information..."):
        # Embed the query
        query_embedding = st.session_state.embeddings.embed_query(query)
        
        # Search the index
        k = min(3, len(st.session_state.chunks))  # Return top 3 results or fewer if we have fewer chunks
        distances, indices = st.session_state.faiss_index.search(
            np.array([query_embedding]).astype('float32'), k
        )
        
        # Get the relevant chunks
        relevant_chunks = [st.session_state.chunks[i] for i in indices[0]]
        
        # Display results
        st.subheader("Search Results")
        for i, (chunk, distance) in enumerate(zip(relevant_chunks, distances[0])):
            st.markdown(f"**Result {i+1}** (Distance: {distance:.4f})")
            st.text_area("", chunk, height=150)
            
        # In a real app, here we would pass these chunks to an LLM to get a coherent answer
        st.info("In a complete implementation, these chunks would be passed to an LLM like GPT-3.5 or GPT-4 to generate a coherent answer to your question.")

# Display parsed documents list
st.sidebar.header("Parsed Documents")
for i, doc in enumerate(st.session_state.parsed_documents):
    st.sidebar.write(f"{i+1}. {doc['name']} ({doc['type']}) - {doc['size']} chars")

# Provide additional information about embedding providers
with st.expander("About Embedding Providers"):
    st.write("""
    ### Amazon Bedrock
    Best for production use with AWS infrastructure. Offers high-quality embedding models like Titan and Cohere.
    
    ### OpenAI
    Good for development and testing. Provides high-quality embeddings but may have higher latency if you're already on AWS.
    
    ### WAS VM (NeMo)
    Use this if you need to keep everything on-premises or have specialized models running on your WAS VM.
    
    ### Mistral API
    Good option for multilingual content with competitive pricing.
    
    For production use with AWS infrastructure, **Amazon Bedrock** is generally recommended for:
    - Seamless AWS integration
    - High-quality embedding models
    - Pay-per-use pricing
    - Automatic scaling
    """)

# Add app information
with st.expander("About this app"):
    st.write("""
    This app demonstrates document parsing, smart text chunking, and vector search capabilities:
    
    1. **Document Parsing**: Extracts text from PDFs, PPTX files, HTML, and other formats
    2. **Smart Text Chunking**: Breaks text into semantically meaningful chunks respecting boundaries like paragraphs and sentences
    3. **Vector Search**: Creates embeddings for chunks and allows searching for relevant information
    
    To use this app:
    1. Select an embedding provider and configure required settings
    2. Upload a document or enter a URL
    3. Click "Parse Document"
    4. Click "Process & Create Embeddings"
    5. Ask questions about the document
    """)